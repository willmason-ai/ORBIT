"""
PPTX extraction — structured text, tables, and RGB fill colors.

The output is fed directly to the Claude agent, so the structure here
is deliberately verbose: each shape carries its raw RGB plus a
best-guess RAG hint if the fill is close enough to a reference color.
"""
from __future__ import annotations

import io
import logging
from typing import Any

from pptx import Presentation

log = logging.getLogger(__name__)

# RAG color reference values (RGB tuples)
RAG_COLORS: dict[str, list[tuple[int, int, int]]] = {
    "RED":   [(255, 0, 0), (192, 0, 0), (255, 0, 16), (196, 15, 38)],
    "AMBER": [(255, 192, 0), (255, 255, 0), (237, 125, 49), (255, 165, 0)],
    "GREEN": [(0, 176, 80), (112, 173, 71), (0, 128, 0), (0, 176, 240)],
}


def _color_distance(c1: tuple[int, int, int], c2: tuple[int, int, int]) -> float:
    return sum((a - b) ** 2 for a, b in zip(c1, c2)) ** 0.5


def _rgb_to_rag(rgb: tuple[int, int, int], threshold: float = 70.0) -> str | None:
    best: str | None = None
    best_dist = float("inf")
    for status, colors in RAG_COLORS.items():
        for ref in colors:
            d = _color_distance(rgb, ref)
            if d < best_dist:
                best_dist, best = d, status
    return best if best_dist <= threshold else None


def extract_pptx(file_bytes: bytes) -> dict[str, Any]:
    """
    Extract all text, table data, and color metadata from a PPTX file.
    Returns a dict ready to be passed to the Claude agent.
    """
    prs = Presentation(io.BytesIO(file_bytes))
    slides_out: list[dict[str, Any]] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data: dict[str, Any] = {
            "slide_number": slide_num,
            "slide_title": None,
            "shapes": [],
            "tables": [],
        }

        try:
            if slide.shapes.title and slide.shapes.title.text:
                slide_data["slide_title"] = slide.shapes.title.text.strip()
        except Exception:
            pass

        for shape in slide.shapes:
            rgb: tuple[int, int, int] | None = None
            rag_hint: str | None = None
            try:
                fill = shape.fill
                if fill.type == 1:  # MSO_FILL.SOLID
                    rgb_obj = fill.fore_color.rgb
                    rgb = (rgb_obj[0], rgb_obj[1], rgb_obj[2])
                    rag_hint = _rgb_to_rag(rgb)
            except Exception:
                pass

            text = ""
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(
                    p.text for p in shape.text_frame.paragraphs if p.text and p.text.strip()
                )

            if getattr(shape, "has_table", False):
                table_data: list[list[str]] = []
                for row in shape.table.rows:
                    table_data.append([cell.text.strip() for cell in row.cells])
                slide_data["tables"].append({
                    "shape_name": shape.name,
                    "rows": table_data,
                })

            if text or rag_hint:
                slide_data["shapes"].append({
                    "shape_name": shape.name,
                    "text": text,
                    "fill_rgb": list(rgb) if rgb else None,
                    "rag_hint": rag_hint,
                })

        slides_out.append(slide_data)

    return {
        "slide_count": len(prs.slides),
        "slides": slides_out,
    }
